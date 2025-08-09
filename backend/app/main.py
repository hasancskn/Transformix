from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import tempfile
import shutil
import subprocess
from pathlib import Path
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from PIL import Image
import io
import os
import logging
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import mm
from pptx import Presentation
from pptx.util import Inches
import pdfplumber
import base64

logger = logging.getLogger("transformix")
logging.basicConfig(level=logging.INFO)

app = FastAPI(title="Transformix API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class Capability(BaseModel):
    name: str
    from_type: str
    to_type: str
    endpoint: str


CAPABILITIES: List[Capability] = [
    Capability(name="Word to PDF", from_type="docx|doc", to_type="pdf", endpoint="/convert/word-to-pdf"),
    Capability(name="PDF to Word", from_type="pdf", to_type="docx", endpoint="/convert/pdf-to-word"),
    Capability(name="PDF to JPG", from_type="pdf", to_type="jpg", endpoint="/convert/pdf-to-jpg"),
    Capability(name="JPG to PDF", from_type="jpg|jpeg|png", to_type="pdf", endpoint="/convert/jpg-to-pdf"),
    Capability(name="Compress PDF", from_type="pdf", to_type="pdf", endpoint="/compress/pdf"),
    Capability(name="Merge PDF", from_type="pdf[]", to_type="pdf", endpoint="/pdf/merge"),
    Capability(name="Split PDF", from_type="pdf", to_type="pdf[]", endpoint="/pdf/split"),
    Capability(name="Rotate PDF", from_type="pdf", to_type="pdf", endpoint="/pdf/rotate"),
    Capability(name="Protect PDF", from_type="pdf", to_type="pdf", endpoint="/pdf/protect"),
    Capability(name="HTML to PDF", from_type="html|url", to_type="pdf", endpoint="/convert/html-to-pdf"),
    # Newly added
    Capability(name="PowerPoint to PDF", from_type="ppt|pptx", to_type="pdf", endpoint="/convert/ppt-to-pdf"),
    Capability(name="Excel to PDF", from_type="xls|xlsx", to_type="pdf", endpoint="/convert/excel-to-pdf"),
    Capability(name="Unlock PDF", from_type="pdf", to_type="pdf", endpoint="/pdf/unlock"),
    Capability(name="Watermark", from_type="pdf", to_type="pdf", endpoint="/pdf/watermark"),
    Capability(name="Page numbers", from_type="pdf", to_type="pdf", endpoint="/pdf/page-numbers"),
    Capability(name="Delete pages", from_type="pdf", to_type="pdf", endpoint="/pdf/delete-pages"),
    Capability(name="Reorder pages", from_type="pdf", to_type="pdf", endpoint="/pdf/reorder"),
    Capability(name="Images to PDF", from_type="jpg|jpeg|png", to_type="pdf", endpoint="/convert/images-to-pdf"),
    Capability(name="PDF to PowerPoint", from_type="pdf", to_type="pptx", endpoint="/convert/pdf-to-pptx"),
    Capability(name="PDF to Excel", from_type="pdf", to_type="xlsx", endpoint="/convert/pdf-to-excel"),
]


@app.get("/")
async def list_capabilities():
    return JSONResponse([c.model_dump() for c in CAPABILITIES])


# Helpers

def _tmp_path(suffix: str) -> Path:
    return Path(tempfile.mkstemp(suffix=suffix)[1])


def _stream_file(path: Path, media_type: str, filename: Optional[str] = None):
    # Read file into memory before TemporaryDirectory cleanup to avoid broken pipe
    data = path.read_bytes()
    headers = {
        "Content-Disposition": f"attachment; filename={filename or path.name}",
    }
    return Response(content=data, media_type=media_type, headers=headers)


# Conversions

@app.post("/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    # Use LibreOffice headless conversion
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        output_dir = Path(tmpdir)
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            result = subprocess.run([
                "soffice", "--headless", "--convert-to", "pdf:writer_pdf_Export", "--outdir", str(output_dir), str(input_path)
            ], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"Conversion failed: {e.stderr.decode(errors='ignore')[:500]}")
        # Expected path by stem
        out_pdf = output_dir / (input_path.stem + ".pdf")
        if not out_pdf.exists():
            # Fallback: pick any single generated PDF
            pdfs = list(output_dir.glob("*.pdf"))
            if len(pdfs) == 1:
                out_pdf = pdfs[0]
        if not out_pdf.exists():
            stdcombined = (result.stdout or b'') + b"\n" + (result.stderr or b'')
            raise HTTPException(500, detail=f"Converted file not found. Logs: {stdcombined.decode(errors='ignore')[:500]}")
        return _stream_file(out_pdf, "application/pdf", filename=out_pdf.name)


@app.post("/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    # Convert PDF to DOCX using pdf2docx (no LibreOffice needed)
    from pdf2docx import Converter
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        output_path = Path(tmpdir) / (Path(file.filename).stem + ".docx")
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            cv = Converter(str(input_path))
            cv.convert(str(output_path), start=0, end=None)
            cv.close()
        except Exception as e:
            raise HTTPException(400, detail=f"pdf2docx failed: {e}")
        if not output_path.exists():
            raise HTTPException(500, detail="Converted file not found")
        return _stream_file(output_path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=output_path.name)


@app.post("/convert/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    # Convert each page to JPG using pdftoppm (Poppler)
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            # Produces files like page-1.jpg, page-2.jpg ...
            subprocess.run([
                "pdftoppm", "-jpeg", "-r", "150", str(input_path), str(Path(tmpdir) / "page")
            ], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"pdftoppm failed: {e.stderr.decode(errors='ignore')[:500]}")
        # Zip results
        import zipfile
        zip_path = Path(tmpdir) / f"{input_path.stem}.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for img in sorted(Path(tmpdir).glob("page-*.jpg")):
                zf.write(img, arcname=img.name)
        return _stream_file(zip_path, "application/zip", filename=zip_path.name)


@app.post("/convert/jpg-to-pdf")
async def jpg_to_pdf(file: UploadFile = File(...)):
    image_bytes = await file.read()
    image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    out_path = _tmp_path(".pdf")
    image.save(out_path, "PDF")
    return _stream_file(out_path, "application/pdf", filename=f"{Path(file.filename).stem}.pdf")


@app.post("/compress/pdf")
async def compress_pdf(file: UploadFile = File(...), quality: int = Form(85)):
    # Use ghostscript for compression
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        output_path = Path(tmpdir) / f"compressed_{Path(file.filename).stem}.pdf"
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            subprocess.run([
                "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/ebook", f"-dJPEGQ={quality}", "-dNOPAUSE", "-dQUIET", "-dBATCH",
                f"-sOutputFile={output_path}", str(input_path)
            ], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"Compression failed: {e.stderr.decode(errors='ignore')[:500]}")
        return _stream_file(output_path, "application/pdf", filename=output_path.name)


@app.post("/pdf/merge")
async def merge_pdf(files: List[UploadFile] = File(...)):
    logger.info("/pdf/merge called with %s files", len(files))
    merger = PdfMerger()
    with tempfile.TemporaryDirectory() as tmpdir:
        for f in files:
            path = Path(tmpdir) / f.filename
            with open(path, "wb") as out:
                shutil.copyfileobj(f.file, out)
            merger.append(str(path))
        merged_path = Path(tmpdir) / "merged.pdf"
        with open(merged_path, "wb") as out:
            merger.write(out)
        merger.close()
        logger.info("/pdf/merge produced %s bytes", merged_path.stat().st_size)
        return _stream_file(merged_path, "application/pdf", filename="merged.pdf")


@app.post("/pdf/split")
async def split_pdf(file: UploadFile = File(...), from_page: int = Form(1), to_page: Optional[int] = Form(None)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        if to_page is None:
            to_page = len(reader.pages)
        writer = PdfWriter()
        for i in range(from_page - 1, to_page):
            writer.add_page(reader.pages[i])
        output = Path(tmpdir) / f"split_{from_page}_{to_page}.pdf"
        with open(output, "wb") as f:
            writer.write(f)
        return _stream_file(output, "application/pdf", filename=output.name)


@app.post("/pdf/rotate")
async def rotate_pdf(file: UploadFile = File(...), degrees: int = Form(90)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        for page in reader.pages:
            page.rotate(degrees)
            writer.add_page(page)
        rotated = Path(tmpdir) / f"rotated_{degrees}.pdf"
        with open(rotated, "wb") as f:
            writer.write(f)
        return _stream_file(rotated, "application/pdf", filename=rotated.name)


@app.post("/pdf/protect")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        protected = Path(tmpdir) / "protected.pdf"
        with open(protected, "wb") as f:
            writer.write(f)
        return _stream_file(protected, "application/pdf", filename=protected.name)


@app.post("/convert/html-to-pdf")
async def html_to_pdf(html: Optional[str] = Form(None), url: Optional[str] = Form(None)):
    logger.info("/convert/html-to-pdf called: url? %s, html? %s", bool(url), bool(html))
    if not html and not url:
        raise HTTPException(400, detail="Provide html or url")
    with tempfile.TemporaryDirectory() as tmpdir:
        output = Path(tmpdir) / "page.pdf"
        common_args = ["--enable-local-file-access", "--encoding", "utf-8", "--quiet", "--custom-header", "User-Agent", "Mozilla/5.0 Transformix"]
        if url:
            try:
                subprocess.run(["wkhtmltopdf", *common_args, url, str(output)], check=True, capture_output=True)
            except subprocess.CalledProcessError as e:
                msg = e.stderr.decode(errors='ignore') or e.stdout.decode(errors='ignore')
                raise HTTPException(400, detail=f"wkhtmltopdf failed: {msg[:500]}")
        else:
            html_path = Path(tmpdir) / "index.html"
            html_path.write_text(html or "", encoding="utf-8")
            try:
                subprocess.run(["wkhtmltopdf", *common_args, str(html_path), str(output)], check=True, capture_output=True)
            except subprocess.CalledProcessError as e:
                msg = e.stderr.decode(errors='ignore') or e.stdout.decode(errors='ignore')
                raise HTTPException(400, detail=f"wkhtmltopdf failed: {msg[:500]}")
        if not output.exists() or output.stat().st_size == 0:
            raise HTTPException(400, detail="wkhtmltopdf produced empty output")
        logger.info("/convert/html-to-pdf produced %s bytes", output.stat().st_size)
        return _stream_file(output, "application/pdf", filename="page.pdf")


@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        output_dir = Path(tmpdir)
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(input_path)], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"Conversion failed: {e.stderr.decode(errors='ignore')[:500]}")
        out_pdf = output_dir / (input_path.stem + ".pdf")
        if not out_pdf.exists():
            pdfs = list(output_dir.glob("*.pdf"))
            if len(pdfs) == 1:
                out_pdf = pdfs[0]
        if not out_pdf.exists():
            raise HTTPException(500, detail="Converted file not found")
        return _stream_file(out_pdf, "application/pdf", filename=out_pdf.name)


@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        output_dir = Path(tmpdir)
        with open(input_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        try:
            subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(input_path)], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"Conversion failed: {e.stderr.decode(errors='ignore')[:500]}")
        out_pdf = output_dir / (input_path.stem + ".pdf")
        if not out_pdf.exists():
            pdfs = list(output_dir.glob("*.pdf"))
            if len(pdfs) == 1:
                out_pdf = pdfs[0]
        if not out_pdf.exists():
            raise HTTPException(500, detail="Converted file not found")
        return _stream_file(out_pdf, "application/pdf", filename=out_pdf.name)


@app.post("/pdf/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        try:
            output = Path(tmpdir) / "unlocked.pdf"
            subprocess.run(["qpdf", f"--password={password}", "--decrypt", str(input_path), str(output)], check=True, capture_output=True)
        except FileNotFoundError:
            raise HTTPException(500, detail="qpdf not installed in image")
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"qpdf failed: {e.stderr.decode(errors='ignore')[:500]}")
        return _stream_file(output, "application/pdf", filename="unlocked.pdf")


@app.post("/pdf/watermark")
async def watermark_pdf(file: UploadFile = File(...), text: Optional[str] = Form(None), image: Optional[UploadFile] = File(None), opacity: float = Form(0.2), size: int = Form(48)):
    if not text and not image:
        raise HTTPException(400, detail="Provide text or image")
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        overlay_path = Path(tmpdir) / "overlay.pdf"
        # Build overlay one-page PDF with reportlab
        c = canvas.Canvas(str(overlay_path), pagesize=letter)
        width, height = letter
        c.saveState()
        c.setFillAlpha(opacity)
        if text:
            c.setFont("Helvetica-Bold", size)
            c.translate(width/2, height/2)
            c.rotate(45)
            c.drawCentredString(0, 0, text)
        if image:
            img_path = Path(tmpdir) / image.filename
            with open(img_path, "wb") as imf:
                shutil.copyfileobj(image.file, imf)
            c.drawImage(str(img_path), width/4, height/4, width=width/2, preserveAspectRatio=True, mask='auto')
        c.restoreState()
        c.showPage()
        c.save()
        # Apply overlay to every page
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        overlay = PdfReader(str(overlay_path)).pages[0]
        for page in reader.pages:
            page.merge_page(overlay)
            writer.add_page(page)
        out_path = Path(tmpdir) / "watermarked.pdf"
        with open(out_path, "wb") as f:
            writer.write(f)
        return _stream_file(out_path, "application/pdf", filename="watermarked.pdf")


@app.post("/pdf/page-numbers")
async def page_numbers(file: UploadFile = File(...), start: int = Form(1), format: str = Form("{n}"), position: str = Form("bottom-right"), size: int = Form(10)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        writer = PdfWriter()
        for idx, page in enumerate(reader.pages, start=start):
            number_text = format.replace("{n}", str(idx))
            # Create small overlay per page
            w = float(page.mediabox.width)
            h = float(page.mediabox.height)
            overlay_path = Path(tmpdir) / f"n_{idx}.pdf"
            c = canvas.Canvas(str(overlay_path), pagesize=(w, h))
            c.setFont("Helvetica", size)
            x, y = w-40, 20
            if position == "bottom-left":
                x, y = 20, 20
            elif position == "top-right":
                x, y = w-40, h-30
            elif position == "top-left":
                x, y = 20, h-30
            c.drawString(x, y, number_text)
            c.save()
            overlay_page = PdfReader(str(overlay_path)).pages[0]
            page.merge_page(overlay_page)
            writer.add_page(page)
        out_path = Path(tmpdir) / "numbered.pdf"
        with open(out_path, "wb") as f:
            writer.write(f)
        return _stream_file(out_path, "application/pdf", filename="numbered.pdf")


@app.post("/pdf/delete-pages")
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    # pages: e.g. "1,3,5-7"
    def parse_pages(s: str, total: int):
        res = set()
        for part in s.split(','):
            part = part.strip()
            if '-' in part:
                a, b = part.split('-')
                for i in range(int(a), int(b) + 1):
                    if 1 <= i <= total:
                        res.add(i)
            else:
                i = int(part)
                if 1 <= i <= total:
                    res.add(i)
        return sorted(res)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        total = len(reader.pages)
        to_delete = set(parse_pages(pages, total))
        writer = PdfWriter()
        for i, p in enumerate(reader.pages, start=1):
            if i not in to_delete:
                writer.add_page(p)
        out_path = Path(tmpdir) / "deleted.pdf"
        with open(out_path, "wb") as f:
            writer.write(f)
        return _stream_file(out_path, "application/pdf", filename="deleted.pdf")


@app.post("/pdf/reorder")
async def reorder_pages(file: UploadFile = File(...), order: str = Form(...)):
    # order: comma-separated indices e.g. "3,1,2"
    new_order = [int(x.strip()) for x in order.split(',') if x.strip()]
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        reader = PdfReader(str(input_path))
        if sorted(new_order) != list(range(1, len(reader.pages)+1)):
            raise HTTPException(400, detail="Order must include each page exactly once")
        writer = PdfWriter()
        for idx in new_order:
            writer.add_page(reader.pages[idx-1])
        out_path = Path(tmpdir) / "reordered.pdf"
        with open(out_path, "wb") as f:
            writer.write(f)
        return _stream_file(out_path, "application/pdf", filename="reordered.pdf")


@app.post("/convert/images-to-pdf")
async def images_to_pdf(files: List[UploadFile] = File(...)):
    images: List[Image.Image] = []
    for f in files:
        img = Image.open(io.BytesIO(await f.read())).convert("RGB")
        images.append(img)
    if not images:
        raise HTTPException(400, detail="No images")
    out_path = _tmp_path(".pdf")
    first, rest = images[0], images[1:]
    first.save(out_path, save_all=True, append_images=rest, format="PDF")
    return _stream_file(out_path, "application/pdf", filename="images.pdf")


@app.post("/convert/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...)):
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        # Render each page to JPG and insert as a slide background
        images_dir = Path(tmpdir) / "imgs"
        images_dir.mkdir(parents=True, exist_ok=True)
        try:
            subprocess.run(["pdftoppm", "-jpeg", "-r", "150", str(input_path), str(images_dir / "page")], check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            raise HTTPException(400, detail=f"pdftoppm failed: {e.stderr.decode(errors='ignore')[:500]}")
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for img in sorted(images_dir.glob("page-*.jpg")):
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(img), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        out_path = Path(tmpdir) / (Path(file.filename).stem + ".pptx")
        prs.save(out_path)
        return _stream_file(out_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=out_path.name)


@app.post("/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    # Simple baseline: extract text lines and write per page to rows in a sheet
    from openpyxl import Workbook
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / file.filename
        with open(input_path, "wb") as out:
            shutil.copyfileobj(file.file, out)
        wb = Workbook()
        ws = wb.active
        ws.title = "Extract"
        row = 1
        try:
            with pdfplumber.open(str(input_path)) as pdf:
                for pageno, page in enumerate(pdf.pages, start=1):
                    ws.cell(row=row, column=1, value=f"Page {pageno}")
                    row += 1
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        ws.cell(row=row, column=1, value=line)
                        row += 1
                    row += 1
        except Exception as e:
            raise HTTPException(400, detail=f"pdfplumber failed: {e}")
        out_path = Path(tmpdir) / (Path(file.filename).stem + ".xlsx")
        wb.save(out_path)
        return _stream_file(out_path, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename=out_path.name) 